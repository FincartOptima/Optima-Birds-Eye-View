"""Overall Portfolio Performance PPTX — starts from the firm's actual client
presentation template and only touches slide 10 ("Optima Actual Performance
Against Benchmark"), replacing its static example images with live tables
built from today's data. Every other slide passes through untouched. A final
disclaimer slide (matching the PDF's) is appended.
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from create_client_factsheet_report import _fmt_inr, _OVERALL_DISCLAIMER_SECTIONS

TEMPLATE_PATH = Path(__file__).resolve().parent / "pptx_templates" / "Overall_Portfolio_Template.pptx"

_NAVY = RGBColor(0x14, 0x36, 0x5C)
_DARK_NAVY = RGBColor(0x0D, 0x24, 0x40)
_HEADER_BLUE = RGBColor(0x1A, 0x3C, 0x8F)
_GOLD = RGBColor(0xC5, 0x92, 0x2E)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_BLUE = RGBColor(0xE8, 0xEE, 0xF6)
_CREAM = RGBColor(0xF3, 0xE9, 0xCE)
_TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
_TEXT_MED = RGBColor(0x5A, 0x6C, 0x7E)
_GREEN = RGBColor(0x1B, 0x7A, 0x2F)
_RED = RGBColor(0xC0, 0x39, 0x2B)


def _set_cell(cell, text, *, bold=False, size=10, color=_TEXT_DARK, fill=None, align=PP_ALIGN.LEFT):
    cell.text = str(text)
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18000)
    cell.margin_bottom = Emu(18000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _return_color(val):
    if val is None:
        return _TEXT_MED
    return _GREEN if val >= 0 else _RED


def _fmt_pct(val):
    if val is None:
        return "N/A"
    return f"{'+' if val >= 0 else ''}{val:.2f}%"


def _remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def _add_label(slide, text, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = _NAVY
    return tb


_ROW_H = Inches(0.32)


def _build_allocation_table(slide, asset_allocation, left, top, width, row_h=None, font_size=10.5):
    row_h = row_h or _ROW_H
    rows = len(asset_allocation) + 2  # header + categories + total
    height = row_h * rows
    graphic_frame = slide.shapes.add_table(rows, 2, left, top, width, height)
    table = graphic_frame.table
    table.columns[0].width = int(width * 0.65)
    table.columns[1].width = int(width * 0.35)

    _set_cell(table.cell(0, 0), "Category", bold=True, size=font_size + 0.5, color=_WHITE, fill=_NAVY)
    _set_cell(table.cell(0, 1), "Weight", bold=True, size=font_size + 0.5, color=_WHITE, fill=_NAVY, align=PP_ALIGN.RIGHT)

    for i, a in enumerate(asset_allocation, start=1):
        bg = _LIGHT_BLUE if i % 2 == 1 else _WHITE
        _set_cell(table.cell(i, 0), a["category"], size=font_size, fill=bg)
        _set_cell(table.cell(i, 1), f"{a['pct']:.2f}%", size=font_size, fill=bg, align=PP_ALIGN.RIGHT)

    last = rows - 1
    _set_cell(table.cell(last, 0), "TOTAL", bold=True, size=font_size, fill=_CREAM)
    _set_cell(table.cell(last, 1), "100.00%", bold=True, size=font_size, fill=_CREAM, align=PP_ALIGN.RIGHT)
    return graphic_frame


def _build_top_holdings_table(slide, top_holdings, left, top, width, row_h=None, font_size=10.5):
    row_h = row_h or _ROW_H
    rows = len(top_holdings) + 1
    height = row_h * rows
    graphic_frame = slide.shapes.add_table(rows, 3, left, top, width, height)
    table = graphic_frame.table
    table.columns[0].width = int(width * 0.08)
    table.columns[1].width = int(width * 0.68)
    table.columns[2].width = int(width * 0.24)

    _set_cell(table.cell(0, 0), "Rank", bold=True, size=font_size + 0.5, color=_WHITE, fill=_NAVY)
    _set_cell(table.cell(0, 1), "Security", bold=True, size=font_size + 0.5, color=_WHITE, fill=_NAVY)
    _set_cell(table.cell(0, 2), "% Assets", bold=True, size=font_size + 0.5, color=_WHITE, fill=_NAVY, align=PP_ALIGN.RIGHT)

    for i, f in enumerate(top_holdings, start=1):
        bg = _LIGHT_BLUE if i % 2 == 1 else _WHITE
        _set_cell(table.cell(i, 0), str(i), size=font_size, fill=bg)
        _set_cell(table.cell(i, 1), f["scheme"], size=font_size, fill=bg)
        _set_cell(table.cell(i, 2), f"{f['pct']:.2f}%", size=font_size, fill=bg, align=PP_ALIGN.RIGHT)
    return graphic_frame


def _build_performance_table(slide, performance, left, top, width, row_h=None, font_size=11, client=None):
    row_h = row_h or Inches(0.5)
    period_labels = performance["period_labels"]
    filtered_benchmarks = [b for b in performance["benchmarks"] if b["name"] in ("BSE 500", "Nifty 50")]

    rows = 1 + 1 + len(filtered_benchmarks)  # header + portfolio row + filtered benchmarks
    cols = 1 + len(period_labels)
    height = row_h * rows
    graphic_frame = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = graphic_frame.table
    table.columns[0].width = int(width * 0.4)
    for i in range(1, cols):
        table.columns[i].width = int(width * 0.6 / len(period_labels))

    _set_cell(table.cell(0, 0), "Portfolio / Benchmark", bold=True, size=font_size + 0.5, color=_WHITE, fill=_DARK_NAVY)
    for i, label in enumerate(period_labels, start=1):
        _set_cell(table.cell(0, i), label, bold=True, size=font_size + 0.5, color=_WHITE, fill=_DARK_NAVY, align=PP_ALIGN.RIGHT)

    # Short label — the header/title already says whose portfolio this is
    _set_cell(table.cell(1, 0), "Portfolio" if client else "Overall Portfolio", bold=True, size=font_size, fill=_CREAM)
    for i, label in enumerate(period_labels, start=1):
        val = performance["portfolio"].get(label)
        _set_cell(table.cell(1, i), _fmt_pct(val), bold=True, size=font_size, fill=_CREAM, color=_return_color(val), align=PP_ALIGN.RIGHT)

    for r, bench in enumerate(filtered_benchmarks, start=2):
        bg = _LIGHT_BLUE if r % 2 == 0 else _WHITE
        _set_cell(table.cell(r, 0), bench["name"], size=font_size, fill=bg)
        for i, label in enumerate(period_labels, start=1):
            val = bench["returns"].get(label)
            _set_cell(table.cell(r, i), _fmt_pct(val), size=font_size, fill=bg, color=_return_color(val), align=PP_ALIGN.RIGHT)
    return graphic_frame


def _rebuild_slide_10(slide, data):
    """Remove the 3 labels + 4 example-data pictures, keep the background /
    header bar / title / logos, then add live tables sized to actually use
    the slide's vertical space instead of leaving one side mostly blank:
    left column stacks Broad Allocation + Performance (both compact, few
    rows each); right column gives Top 5 Holdings the full column height
    with generously large rows, since it's the only thing there and has
    the longest text (fund names) that benefits most from the room. The
    old performance chart is dropped, not replaced."""
    shapes_by_name = {s.name: s for s in slide.shapes}
    for name in ("TextBox 9", "TextBox 11", "TextBox 13", "Picture 19", "Picture 5", "Picture 6", "Picture 3"):
        if name in shapes_by_name:
            _remove_shape(shapes_by_name[name])

    label_h = Inches(0.4)
    label_gap = Inches(0.1)

    # ---- Left column: Broad Allocation (top) + Performance (bottom) ----
    left_x, left_w = Inches(0.3), Inches(6.0)

    alloc_label_top = Inches(1.0)
    alloc_table_top = alloc_label_top + label_h + label_gap
    _add_label(slide, "Broad Allocation", left_x, alloc_label_top, Inches(3.5), label_h)
    alloc_row_h = Inches(0.35)
    _build_allocation_table(
        slide, data["asset_allocation"],
        left_x, alloc_table_top, left_w,
        row_h=alloc_row_h, font_size=12,
    )

    n_alloc_rows = len(data["asset_allocation"]) + 2
    alloc_table_bottom = alloc_table_top + alloc_row_h * n_alloc_rows

    inception = data.get("inception_date", "")
    perf_label_top = alloc_table_bottom + Inches(0.35)
    perf_table_top = perf_label_top + label_h + label_gap
    _add_label(slide, f"Performance (Inception: {inception})", left_x, perf_label_top, Inches(5.5), label_h)
    _build_performance_table(
        slide, data["performance"],
        left_x, perf_table_top, left_w,
        row_h=Inches(0.5), font_size=12, client=data.get("client"),
    )

    # ---- Right column: Top 5 Holdings alone, given the full column height ----
    right_x, right_w = Inches(6.55), Inches(6.5)
    top5_label_top = Inches(1.0)
    top5_table_top = top5_label_top + label_h + label_gap
    _add_label(slide, "Top 5 Holdings", right_x, top5_label_top, Inches(4), label_h)
    _build_top_holdings_table(
        slide, data["top_holdings"],
        right_x, top5_table_top, right_w,
        row_h=Inches(0.7), font_size=13,
    )


def _copy_shape(source_shape, target_slide):
    new_el = copy.deepcopy(source_shape._element)
    target_slide.shapes._spTree.append(new_el)


def _add_disclaimer_slide(prs, reference_slide):
    layout = reference_slide.slide_layout
    slide = prs.slides.add_slide(layout)

    ref_by_name = {s.name: s for s in reference_slide.shapes}
    for name in ("Shape 1", "Shape 137", "Image 16", "Image 17"):
        if name in ref_by_name:
            _copy_shape(ref_by_name[name], slide)

    title = slide.shapes.add_textbox(Inches(0.52), Inches(0.17), Inches(11.13), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Disclaimer & Statutory Disclosures"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = _WHITE

    # Two columns, 3 sections each — halves the vertical text load per column
    # so there's a real safety margin against overflowing the slide, and
    # reads better on a slide this wide anyway.
    midpoint = (len(_OVERALL_DISCLAIMER_SECTIONS) + 1) // 2
    columns = [
        (Inches(0.4), _OVERALL_DISCLAIMER_SECTIONS[:midpoint]),
        (Inches(6.75), _OVERALL_DISCLAIMER_SECTIONS[midpoint:]),
    ]
    for col_left, sections in columns:
        body = slide.shapes.add_textbox(col_left, Inches(1.0), Inches(6.15), Inches(6.3))
        tf = body.text_frame
        tf.word_wrap = True
        first = True
        for heading, text in sections:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = heading
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = _NAVY

            p2 = tf.add_paragraph()
            p2.space_after = Pt(10)
            run2 = p2.add_run()
            run2.text = text
            run2.font.size = Pt(10)
            run2.font.color.rgb = _TEXT_MED


def generate_overall_pptx(data: dict, output_path: Path) -> None:
    prs = Presentation(str(TEMPLATE_PATH))
    slide10 = prs.slides[9]
    _rebuild_slide_10(slide10, data)
    _add_disclaimer_slide(prs, slide10)
    prs.save(str(output_path))
